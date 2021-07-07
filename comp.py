# required: use pip to install python-pptx, pywin32, opencv-python, pytesseract, pil, tesseract, tesseract-ocr, natsort

from pptx import Presentation
import win32com
import win32com.client
import os
import sys
import cv2
import glob
import natsort
import argparse
import pytesseract
from difflib import SequenceMatcher
from datetime import datetime
import argparse

parser = argparse.ArgumentParser()
# add arguments for running
parser.add_argument('-video', '--video', help='Video file name', required=True)
parser.add_argument('-ppt', '--ppt', help='Powerpoint file name', required=True)
args = vars(parser.parse_args())
startTime = datetime.now()

# SET TESSERACT PATH HERE
pytesseract.pytesseract.tesseract_cmd = r'F:\Program Files\Tesseract-OCR\tesseract.exe'

filename = args['ppt']
titles = []
width = 0
height = 0


def get_text(img):
    """
    Get text string read from a image file read by cv2 using pytesseract image
    :img: cv2 image
    :returns: string containing text from the image
    """
    return pytesseract.image_to_string(img)


def similar(a, b):
    """
    Get string similarity from two strings
    :a: str 1 for comparison
    :b: str 2 for comparison
    :returns: decimal similarity of the two strings
    """
    return SequenceMatcher(None, a, b).ratio()


def extractImages(pathIn, pathOut):
    """
    Extract image screenshots from a video file
    :pathIn: path for the video file
    :pathOut: path for storing all the screenshots
    """
    count = 0
    vidcap = cv2.VideoCapture(pathIn)
    success, image = vidcap.read()
    success = True
    width = vidcap.get(cv2.CAP_PROP_FRAME_WIDTH)
    height = vidcap.get(cv2.CAP_PROP_FRAME_HEIGHT)
    print("WxH" + str(width) + str(height))
    prev = None
    mid = int(width) / 2
    # keep reading while there are frames available
    while success:
        vidcap.set(cv2.CAP_PROP_POS_MSEC, (count * 1000))
        success, image = vidcap.read()
        print ('Read a new frame: ', success)
        if (success and count % 5 == 0):
            # if image is dark at the sample point turn it white for better OCR
            if (image[80, int(width - 150), 1] > 127):
                image = (255 - image)
            image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)

            ret, image = cv2.threshold(image, 200, 255, 1)
            cv2.imwrite(pathOut + "\\%d.jpg" % count, image)     # save frame as JPEG file
            prev = image
        count = count + 1


def filename_to_time(name):
    """
    Turns filename (stored in seconds) to proper time format
    :name: string filename
    :returns: string containing time stored in min:sec format
    """
    time_str = name[:len(name) - 4]
    time_int = int(time_str)
    mins = time_int // 60
    seconds = time_int % 60
    time = str(mins) + ":" + str(seconds)
    return time


def get_text_from_image(filename):
    """
    Reads an image file, gets its text and cleans it up
    :filename: file name of image being read
    :returns: cleaned up text of the image that was read
    """
    full_name = r'ss/' + str(filename)
    print(full_name)
    image = cv2.imread(full_name)
    text = get_text(image)
    text = text.replace('\n', ' ').replace('\x0c', '').replace('—', '–')
    return text


def get_inverted_text(filename):
    """
    Reads an image file, reverses colours, gets its text and cleans it up
    Currently unused
    :filename: file name of image being read
    :returns: cleaned up text of the image that was read and reversed
    """
    full_name = r'ss/' + str(filename)
    print(full_name)
    image = cv2.imread(full_name)
    image = (255 - image)
    text = get_text(image)
    text = text.replace('\n', ' ').replace('\x0c', '')
    return text


def clear_files():
    """Cleans all files in the screenshots directory"""
    files = glob.glob('ss/*')
    for f in files:
        os.remove(f)


prs = Presentation(filename)

# find titles by getting title object
for slide in prs.slides:
    title = slide.shapes.title
    if title != None:
        title = slide.shapes.title.text
        title = title.replace("\x0b", "").strip(""" !@#$%^&*)(_-+=}{][:;<,>.?"'/<,""")
        titles.append(title)

# if title finding wasn't successful, use 0th index instead
if titles == []:
    for slide in prs.slides:
        title_shape = slide.shapes[0]  # consider the zeroth indexed shape as the title
        if title_shape.has_text_frame:
            if title_shape.text.strip(""" !@#$%^&*)(_-+=}{][:;<,>.?"'/<,""") not in titles:
                titles.append(title_shape.text.strip(""" !@#$%^&*)(_-+=}{][:;<,>.?"'/<,"""))

titles = list(dict.fromkeys(titles))

titles_length = len(titles)
print(titles)

# cleanup and make new screenshots in the ss folder
clear_files()
if not os.path.exists('ss'):
    os.mkdir('ss')
extractImages(args['video'], r'ss')


files = [f for f in os.listdir('ss') if os.path.isfile(os.path.join('ss', f))]
sorted_files = natsort.natsorted(files)
print(sorted_files)

# remove images that are too similar


def remove_similar():
    """Iterates through images in a folder that are similar to one another"""
    prev = None
    for i in range(0, len(sorted_files)):
        if prev == None:
            prev = sorted_files[i]
        else:
            curr_name = r'ss/' + str(sorted_files[i])
            prev_name = r'ss/' + prev
            a = cv2.imread(curr_name)
            b = cv2.imread(prev_name)
            if (cv2.absdiff(a, b).mean() < 1):  # threshold currently 1
                os.remove(curr_name)
            else:
                prev = sorted_files[i]


remove_similar()
files = [f for f in os.listdir('ss') if os.path.isfile(os.path.join('ss', f))]
sorted_files = natsort.natsorted(files)
print(sorted_files)

images_and_strings = []

for i in sorted_files:
    text = get_text_from_image(i)
    tup = (i, text)
    images_and_strings.append(tup)

result = []
found_list = []

# find matches between titles extracted from ppt and strings extracted from screenshots
for i in titles:
    found = False
    j = 0
    while not found and j < len(images_and_strings):
        if i.casefold() in images_and_strings[j][1].casefold():
            repeat = False
            for g in found_list:
                if g.casefold() in images_and_strings[j][1].casefold():
                    # check for repeats, if it repeats then skip
                    if (images_and_strings[j][1].casefold().index(g.casefold()) < images_and_strings[j][1].casefold().index(i.casefold())):
                        repeat = True
            if not repeat:
                found = True
                res_tup = (i, filename_to_time(images_and_strings[j][0]))
                result.append(res_tup)
                found_list.append(i)
            j += 1
        else:
            j += 1

print(result)

res_path = "result.txt"
if os.path.isfile(res_path):
    os.remove(res_path)

# write results to file
re = open(res_path, 'w')
for i in result:
    re.write(i[0] + ' - ' + i[1] + '\n')
re.close()

print("Time taken: " + str((datetime.now() - startTime)))
